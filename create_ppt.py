# -*- coding: utf-8 -*-
"""
生成 JavaWeb 购物网站答辩 PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 宽屏
prs.slide_height = Inches(7.5)

# ==================== 颜色方案 ====================
PRIMARY = RGBColor(0x09, 0x84, 0xE3)      # 主色蓝
PRIMARY_DARK = RGBColor(0x07, 0x70, 0xC0)
ACCENT = RGBColor(0x6C, 0x5C, 0xE7)        # 紫色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x2D, 0x34, 0x36)
GRAY = RGBColor(0x63, 0x6E, 0x72)
LIGHT_GRAY = RGBColor(0xF5, 0xF6, 0xFA)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE1, 0x70, 0x55)
DARK_BG = RGBColor(0x1E, 0x27, 0x2E)


def add_bg(slide, color):
    """给幻灯片添加纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=BLACK,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_multi_text(slide, left, top, width, height, lines, font_name='Microsoft YaHei'):
    """添加多行文本框，lines 是 [(text, font_size, color, bold, alignment), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text, font_size, color, bold, alignment = line_data
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
    return tf


def add_bottom_bar(slide):
    """底部装饰条"""
    add_shape_bg(slide, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3), PRIMARY)


def add_page_number(slide, num, total):
    """页码"""
    add_text_box(slide, Inches(12.2), Inches(7.22), Inches(1), Inches(0.25),
                 f"{num}/{total}", font_size=10, color=WHITE, alignment=PP_ALIGN.RIGHT)


def add_title_bar(slide, title, subtitle=""):
    """顶部标题栏"""
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.6),
                 title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.65), Inches(11), Inches(0.4),
                     subtitle, font_size=14, color=RGBColor(0xDF, 0xE6, 0xE9))


# ============================================================
# 幻灯片 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_bg(slide, DARK_BG)

# 顶部装饰线
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

# 左侧装饰竖条
add_shape_bg(slide, Inches(1.5), Inches(2.0), Inches(0.06), Inches(3.8), PRIMARY)

# 项目名称
add_text_box(slide, Inches(2.0), Inches(2.0), Inches(10), Inches(0.8),
             "简易购物网站", font_size=48, color=WHITE, bold=True)

# 副标题
add_text_box(slide, Inches(2.0), Inches(2.9), Inches(10), Inches(0.6),
             "ShoppingWeb — 基于 Servlet + JSP + MySQL 的电子商务平台",
             font_size=20, color=RGBColor(0xB2, 0xBE, 0xC3))

# 分隔线
add_shape_bg(slide, Inches(2.0), Inches(3.7), Inches(3), Inches(0.03), PRIMARY)

# 信息区
add_multi_text(slide, Inches(2.0), Inches(4.1), Inches(8), Inches(2.5), [
    ("Java Web 期末课程设计答辩", 22, RGBColor(0xDF, 0xE6, 0xE9), False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    (f"答辩日期：{datetime.date.today().strftime('%Y年%m月%d日')}", 16, GRAY, False, PP_ALIGN.LEFT),
    ("技术栈：Java 8 + Servlet 4.0 + JSP 2.3 + MySQL 8.0 + Maven + Tomcat 9", 14, GRAY, False, PP_ALIGN.LEFT),
])

# 底部
add_shape_bg(slide, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3), PRIMARY)


# ============================================================
# 幻灯片 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "目  录", "CONTENTS")

items = [
    ("01", "项目背景与需求分析", "项目选题背景、目标用户与核心需求"),
    ("02", "技术架构与开发环境", "技术选型、系统架构、开发工具"),
    ("03", "数据库设计", "ER 图、表结构设计、关系说明"),
    ("04", "功能模块详解", "用户端 + 管理端核心功能流程"),
    ("05", "核心代码展示", "关键代码逻辑与实现亮点"),
    ("06", "界面展示", "前端页面效果截图展示"),
    ("07", "项目总结与展望", "收获、不足、改进方向"),
]

for i, (num, title, desc) in enumerate(items):
    y = Inches(1.5) + Inches(0.75) * i
    # 编号圆圈
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.8), y, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY if i < 6 else ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Microsoft YaHei'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(2.6), y + Inches(0.02), Inches(6), Inches(0.3),
                 title, font_size=20, color=BLACK, bold=True)
    add_text_box(slide, Inches(2.6), y + Inches(0.32), Inches(7), Inches(0.2),
                 desc, font_size=12, color=GRAY)

add_bottom_bar(slide)
add_page_number(slide, 2, 8)


# ============================================================
# 幻灯片 3: 项目背景与需求分析
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "项目背景与需求分析", "BACKGROUND & REQUIREMENTS")

# 左侧 - 项目背景
add_shape_bg(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.3), LIGHT_GRAY)
add_text_box(slide, Inches(0.9), Inches(1.7), Inches(5), Inches(0.4),
             "📋 项目背景", font_size=22, color=PRIMARY, bold=True)
add_multi_text(slide, Inches(0.9), Inches(2.2), Inches(5.2), Inches(4.2), [
    ("● 电子商务持续高速发展，网上购物已成为主流消费方式", 14, BLACK, False, PP_ALIGN.LEFT),
    ("● Java Web 技术在企业级开发中广泛应用，是计算机专业核心课程", 14, BLACK, False, PP_ALIGN.LEFT),
    ("● 通过实际项目锻炼 Servlet、JSP、JDBC 等技术的综合运用能力", 14, BLACK, False, PP_ALIGN.LEFT),
    ("● 加深对 MVC 分层架构和前后端协作的理解", 14, BLACK, False, PP_ALIGN.LEFT),
    ("● 积累完整的项目开发流程经验：需求→设计→编码→测试→部署", 14, BLACK, False, PP_ALIGN.LEFT),
])

# 右侧 - 功能需求
add_shape_bg(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), LIGHT_GRAY)
add_text_box(slide, Inches(7.1), Inches(1.7), Inches(5), Inches(0.4),
             "🎯 核心功能需求", font_size=22, color=ACCENT, bold=True)

add_multi_text(slide, Inches(7.1), Inches(2.2), Inches(5.2), Inches(4.5), [
    ("👤 用户端功能", 16, BLACK, True, PP_ALIGN.LEFT),
    ("  • 用户注册与登录（MD5 加密）", 13, BLACK, False, PP_ALIGN.LEFT),
    ("  • 商品浏览（分页 + 分类筛选 + 关键词搜索）", 13, BLACK, False, PP_ALIGN.LEFT),
    ("  • 商品详情查看", 13, BLACK, False, PP_ALIGN.LEFT),
    ("  • 购物车管理（添加 / 修改数量 / 删除）", 13, BLACK, False, PP_ALIGN.LEFT),
    ("  • 订单结算（事务控制：扣库存 + 清空购物车）", 13, BLACK, False, PP_ALIGN.LEFT),
    ("  • 订单查看与详情", 13, BLACK, False, PP_ALIGN.LEFT),
    ("", 6, BLACK, False, PP_ALIGN.LEFT),
    ("🛠️ 管理端功能", 16, BLACK, True, PP_ALIGN.LEFT),
    ("  • 商品管理（增删改查 CRUD）", 13, BLACK, False, PP_ALIGN.LEFT),
    ("  • 订单管理（查看所有订单 + 修改订单状态）", 13, BLACK, False, PP_ALIGN.LEFT),
])

add_bottom_bar(slide)
add_page_number(slide, 3, 8)


# ============================================================
# 幻灯片 4: 技术架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "技术架构与开发环境", "TECHNICAL ARCHITECTURE")

# 左侧 - 技术栈表格
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4),
             "🛠️ 技术栈", font_size=22, color=PRIMARY, bold=True)

tech_data = [
    ("Java 8+", "后端开发语言"),
    ("Servlet 4.0", "Web 控制器，处理 HTTP 请求"),
    ("JSP 2.3 + JSTL 1.2", "视图层模板引擎"),
    ("MySQL 8.0", "关系型数据库"),
    ("JDBC", "数据库连接与操作"),
    ("Maven 3.x", "项目构建与依赖管理"),
    ("Tomcat 9", "Web 应用容器"),
    ("HTML5 + CSS3 + JS", "前端页面与交互"),
]

table_shape = slide.shapes.add_table(len(tech_data) + 1, 2,
                                      Inches(0.6), Inches(2.1), Inches(5.8), Inches(4.5))
table = table_shape.table
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(3.6)

# 表头
for j, header in enumerate(["技术/工具", "说明"]):
    cell = table.cell(0, j)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = PRIMARY
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER

for i, (tech, desc) in enumerate(tech_data):
    for j, val in enumerate([tech, desc]):
        cell = table.cell(i + 1, j)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_GRAY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = BLACK
            p.font.name = 'Microsoft YaHei'
            p.alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT

# 右侧 - 架构图
add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5), Inches(0.4),
             "🏗️ 分层架构", font_size=22, color=ACCENT, bold=True)

# 架构分层图
layers = [
    ("表示层 (View)", "JSP 页面 + CSS + JavaScript\nJSTL 标签库", RGBColor(0x09, 0x84, 0xE3)),
    ("控制层 (Controller)", "Servlet 控制器\n处理请求、验证、转发", RGBColor(0x6C, 0x5C, 0xE7)),
    ("业务逻辑层 (DAO)", "数据访问对象\n封装数据库 CRUD 操作", RGBColor(0x00, 0xB8, 0x94)),
    ("数据层 (Database)", "MySQL 数据库\n6 张数据表 + 外键约束", RGBColor(0xE1, 0x70, 0x55)),
]

for i, (title, desc, color) in enumerate(layers):
    y = Inches(2.1) + Inches(1.2) * i
    # 方块
    box = add_shape_bg(slide, Inches(7.2), y, Inches(5.2), Inches(1.0), color)
    add_text_box(slide, Inches(7.5), y + Inches(0.08), Inches(4.6), Inches(0.3),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, Inches(7.5), y + Inches(0.38), Inches(4.6), Inches(0.55),
                 desc, font_size=11, color=RGBColor(0xDF, 0xE6, 0xE9))
    # 箭头（除了最后一层）
    if i < len(layers) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                        Inches(9.5), y + Inches(1.0), Inches(0.4), Inches(0.2))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()

add_bottom_bar(slide)
add_page_number(slide, 4, 8)


# ============================================================
# 幻灯片 5: 数据库设计
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "数据库设计", "DATABASE DESIGN")

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.4),
             "📊 数据库关系模型（6 张核心表）", font_size=20, color=PRIMARY, bold=True)

# 表结构展示 - 使用3列卡片布局
tables_info = [
    ("users 用户表", PRIMARY, [
        "id (PK, AUTO_INCREMENT)",
        "username (UNIQUE, NOT NULL)",
        "password (MD5加密存储)",
        "email / phone / address",
        "role (user / admin)",
        "create_time",
    ]),
    ("categories 分类表", ACCENT, [
        "id (PK, AUTO_INCREMENT)",
        "name (分类名称)",
        "sort (排序序号)",
        "5个预设分类：电子、服装、\n图书、家居、零食",
    ]),
    ("products 商品表", GREEN, [
        "id (PK, AUTO_INCREMENT)",
        "name / description",
        "price (DECIMAL 10,2)",
        "stock (库存数量)",
        "category_id (FK → categories)",
        "image / create_time",
    ]),
    ("cart_items 购物车表", ORANGE, [
        "id (PK, AUTO_INCREMENT)",
        "user_id (FK → users)",
        "product_id (FK → products)",
        "quantity (数量)",
        "add_time",
    ]),
    ("orders 订单表", RGBColor(0xA2, 0x9B, 0xFE), [
        "id (PK, AUTO_INCREMENT)",
        "order_no (UNIQUE 订单编号)",
        "user_id (FK → users)",
        "total_price",
        "status (pending/paid/shipped/\ncompleted/canceled)",
        "address / create_time",
    ]),
    ("order_items 订单明细表", RGBColor(0xFD, 0x79, 0x89), [
        "id (PK, AUTO_INCREMENT)",
        "order_id (FK → orders)",
        "product_id (FK → products)",
        "product_name (快照)",
        "product_price (下单时价格)",
        "quantity",
    ]),
]

for idx, (name, color, fields) in enumerate(tables_info):
    col = idx % 3
    row = idx // 3
    x = Inches(0.5) + Inches(4.15) * col
    y = Inches(2.0) + Inches(2.6) * row

    # 卡片背景
    add_shape_bg(slide, x, y, Inches(3.85), Inches(2.3), WHITE)
    # 顶部色条
    add_shape_bg(slide, x, y, Inches(3.85), Inches(0.06), color)
    # 表名
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(3.5), Inches(0.3),
                 name, font_size=15, color=color, bold=True)
    # 字段列表
    field_text = "\n".join(f"• {f}" for f in fields)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(3.5), Inches(1.7),
                 field_text, font_size=11, color=BLACK)

# 关键设计说明
add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.25),
             "💡 设计亮点：密码 MD5 加密 | 订单商品名称快照（防止商品修改后历史订单数据变化）| 外键约束保证数据完整性 | 事务保证下单原子性",
             font_size=11, color=GRAY)

add_page_number(slide, 5, 8)


# ============================================================
# 幻灯片 6: 功能模块详解
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "功能模块详解", "FUNCTIONAL MODULES")

# 用户端流程
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.4),
             "👤 用户端核心流程", font_size=20, color=PRIMARY, bold=True)

user_steps = [
    ("1", "注册/登录", "表单验证 → MD5加密 → Session保存", PRIMARY),
    ("2", "浏览商品", "分页查询 + 分类筛选 + 关键词搜索 + 排序", ACCENT),
    ("3", "加入购物车", "库存校验 → 插入购物车表 → AJAX 刷新角标", GREEN),
    ("4", "结算下单", "库存二次校验 → 事务写入 → 扣库存 → 清空购物车", ORANGE),
    ("5", "查看订单", "订单列表 → 订单详情（含商品快照）", RGBColor(0xA2, 0x9B, 0xFE)),
]

for i, (num, title, desc, color) in enumerate(user_steps):
    y = Inches(2.0) + Inches(0.9) * i
    # 序号圆圈
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), y, Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Microsoft YaHei'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(1.5), y + Inches(0.02), Inches(2), Inches(0.25),
                 title, font_size=15, color=BLACK, bold=True)
    add_text_box(slide, Inches(1.5), y + Inches(0.28), Inches(5), Inches(0.2),
                 desc, font_size=11, color=GRAY)

# 管理端流程
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.4),
             "🛠️ 管理端功能", font_size=20, color=ACCENT, bold=True)

admin_modules = [
    ("📦 商品管理", [
        "商品列表（后台分页）",
        "新增商品（名称/价格/库存/分类）",
        "编辑商品信息",
        "删除商品",
    ]),
    ("📋 订单管理", [
        "查看所有用户订单",
        "查看订单详情",
        "修改订单状态（待付款→已付款→已发货→已完成→已取消）",
        "关联显示下单用户信息",
    ]),
]

for i, (title, items) in enumerate(admin_modules):
    y = Inches(2.0) + Inches(2.3) * i
    add_shape_bg(slide, Inches(7.0), y, Inches(5.6), Inches(2.0), LIGHT_GRAY)
    add_text_box(slide, Inches(7.3), y + Inches(0.15), Inches(5), Inches(0.3),
                 title, font_size=16, color=ACCENT, bold=True)
    item_text = "\n".join(f"  ✓ {item}" for item in items)
    add_text_box(slide, Inches(7.3), y + Inches(0.55), Inches(5), Inches(1.4),
                 item_text, font_size=12, color=BLACK)

# 安全控制
add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
             "🔒 安全控制：EncodingFilter（全局编码）| LoginFilter（JSP页面拦截 + 登录后回跳）| AdminFilter（角色权限校验 + 越权拦截）",
             font_size=12, color=RED)

add_bottom_bar(slide)
add_page_number(slide, 6, 8)


# ============================================================
# 幻灯片 7: 核心代码展示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "核心代码展示", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.3),
             "KEY CODE HIGHLIGHTS", font_size=14, color=GRAY)

# 三个代码卡片
code_cards = [
    ("🔐 下单事务控制 (OrderDao.createOrder)",
     GREEN,
     """// 创建订单（事务：4步原子操作）
conn.setAutoCommit(false);  // 开启事务
// 1. 插入订单记录
String orderNo = "ORD" + timestamp + uuid;
INSERT INTO orders (...) VALUES (...);
// 2. 插入订单明细（商品快照）
INSERT INTO order_items (...) VALUES (...);
// 3. 扣减库存（带库存校验）
UPDATE products SET stock = stock - ?
WHERE id = ? AND stock >= ?;
// 4. 清空用户购物车
DELETE FROM cart_items WHERE user_id = ?;
conn.commit();  // 提交事务
// 任何一步失败 → conn.rollback() 全部回滚"""),

    ("🔍 商品搜索 + 分页 (ProductDao.findAll)",
     ACCENT,
     """// 动态 SQL 构建
StringBuilder sql = new StringBuilder(
  "SELECT p.*, c.name FROM products p
   LEFT JOIN categories c ON ... WHERE 1=1");
// 可选筛选条件
if (categoryId != null)
  sql.append(" AND p.category_id = ?");
if (keyword != null)
  sql.append(" AND (p.name LIKE ?
    OR p.description LIKE ?)");
// 分页
sql.append(" ORDER BY p.id DESC
  LIMIT ?, ?");  // LIMIT offset, pageSize"""),

    ("🛡️ 多层过滤器链 (Filter Chain)",
     ORANGE,
     """// EncodingFilter — 全局 UTF-8 编码
request.setCharacterEncoding("UTF-8");
response.setCharacterEncoding("UTF-8");
chain.doFilter(req, resp);

// LoginFilter — 登录拦截
User user = session.getAttribute("user");
if (user == null) {
  session.setAttribute("redirectUrl", uri);
  redirect("/login.jsp"); return;
}

// AdminFilter — 管理员权限
if (!user.isAdmin()) {
  alert("权限不足！"); return;
}"""),
]

for i, (title, color, code) in enumerate(code_cards):
    x = Inches(0.5) + Inches(4.2) * i
    y = Inches(1.5)

    # 顶部色条
    add_shape_bg(slide, x, y, Inches(3.9), Inches(0.06), color)
    # 卡片背景
    add_shape_bg(slide, x, y + Inches(0.06), Inches(3.9), Inches(5.5), RGBColor(0x2C, 0x36, 0x3E))
    # 标题
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.5), Inches(0.3),
                 title, font_size=13, color=color, bold=True)
    # 代码
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(3.5), Inches(4.7),
                 code, font_size=10, color=RGBColor(0xCD, 0xD6, 0xE0))

add_page_number(slide, 7, 8)


# ============================================================
# 幻灯片 8: 界面展示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "界面展示", "USER INTERFACE")

pages_demo = [
    ("🏠 首页", "渐变横幅 Hero Banner\n热门推荐商品网格（4列）\n全局搜索框 + 购物车角标", PRIMARY),
    ("🛍️ 商品列表", "分类胶囊导航栏\n商品卡片网格布局\n分页组件\n关键词搜索高亮", ACCENT),
    ("📦 商品详情", "左右双栏布局\n数量加减控件\n加入购物车按钮\n库存状态实时显示", GREEN),
    ("🛒 购物车", "表格式商品列表\n数量修改/删除操作\n总价实时计算\n结算下单按钮", ORANGE),
    ("📋 订单管理", "订单列表 + 状态标签\n订单详情（商品快照）\n状态流转时间线\n收货地址信息", RGBColor(0xA2, 0x9B, 0xFE)),
    ("⚙️ 管理后台", "管理仪表盘导航\n商品 CRUD 表单\n订单状态修改\n全部订单视图", RGBColor(0xFD, 0x79, 0x89)),
]

for i, (name, desc, color) in enumerate(pages_demo):
    col = i % 3
    row = i // 3
    x = Inches(0.5) + Inches(4.2) * col
    y = Inches(1.5) + Inches(2.9) * row

    # 卡片
    add_shape_bg(slide, x, y, Inches(3.9), Inches(2.5), WHITE)
    # 色条
    add_shape_bg(slide, x, y, Inches(3.9), Inches(0.06), color)
    # 模拟浏览器窗口
    add_shape_bg(slide, x + Inches(0.15), y + Inches(0.3), Inches(3.6), Inches(0.2), LIGHT_GRAY)
    # 浏览器圆点
    for dot_x in [x + Inches(0.25), x + Inches(0.4), x + Inches(0.55)]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, dot_x, y + Inches(0.34), Inches(0.08), Inches(0.08))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GRAY
        dot.line.fill.background()

    # 页面内容模拟区
    add_shape_bg(slide, x + Inches(0.15), y + Inches(0.65), Inches(3.6), Inches(1.2), LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.75), Inches(3.3), Inches(0.3),
                 name, font_size=15, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.1), Inches(3.3), Inches(0.7),
                 desc, font_size=11, color=BLACK)

add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.25),
             "🎨 前端特色：响应式布局（适配PC/平板/手机）| CSS Grid 商品网格 | 渐变色彩设计 | 悬停动效 | 统一组件风格",
             font_size=11, color=GRAY)

add_page_number(slide, 8, 8)


# ============================================================
# 幻灯片 9: 项目总结
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "项目总结与展望", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.3),
             "SUMMARY & FUTURE WORK", font_size=14, color=GRAY)

# 收获
add_shape_bg(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.8), RGBColor(0x2C, 0x36, 0x3E))
add_text_box(slide, Inches(1.0), Inches(1.8), Inches(5), Inches(0.3),
             "✅ 项目收获", font_size=20, color=GREEN, bold=True)
add_multi_text(slide, Inches(1.0), Inches(2.3), Inches(5.2), Inches(1.8), [
    ("✓ 掌握了 Servlet + JSP Web 开发完整流程", 14, RGBColor(0xCD, 0xD6, 0xE0), False, PP_ALIGN.LEFT),
    ("✓ 深入理解 MVC 分层架构设计思想", 14, RGBColor(0xCD, 0xD6, 0xE0), False, PP_ALIGN.LEFT),
    ("✓ 熟练运用 JDBC 进行数据库操作与事务控制", 14, RGBColor(0xCD, 0xD6, 0xE0), False, PP_ALIGN.LEFT),
    ("✓ 学会 Filter 过滤器实现权限校验", 14, RGBColor(0xCD, 0xD6, 0xE0), False, PP_ALIGN.LEFT),
    ("✓ 掌握了前端响应式布局与 CSS 设计技巧", 14, RGBColor(0xCD, 0xD6, 0xE0), False, PP_ALIGN.LEFT),
    ("✓ 积累了 Maven 构建与 Tomcat 部署经验", 14, RGBColor(0xCD, 0xD6, 0xE0), False, PP_ALIGN.LEFT),
])

# 不足
add_shape_bg(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.8), RGBColor(0x2C, 0x36, 0x3E))
add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5), Inches(0.3),
             "⚠️ 现有不足", font_size=20, color=ORANGE, bold=True)
add_multi_text(slide, Inches(7.2), Inches(2.3), Inches(5.2), Inches(1.8), [
    ("✗ 未实现图片上传功能（使用占位图）", 14, RGBColor(0xCD, 0xD6, 0xE0), False, PP_ALIGN.LEFT),
    ("✗ 缺少真正的支付接口集成", 14, RGBColor(0xCD, 0xD6, 0xE0), False, PP_ALIGN.LEFT),
    ("✗ 无用户个人信息管理页面", 14, RGBColor(0xCD, 0xD6, 0xE0), False, PP_ALIGN.LEFT),
    ("✗ 未使用连接池（每次新建连接）", 14, RGBColor(0xCD, 0xD6, 0xE0), False, PP_ALIGN.LEFT),
    ("✗ 缺少单元测试覆盖", 14, RGBColor(0xCD, 0xD6, 0xE0), False, PP_ALIGN.LEFT),
])

# 展望
add_shape_bg(slide, Inches(0.6), Inches(4.7), Inches(12.0), Inches(2.2), RGBColor(0x2C, 0x36, 0x3E))
add_text_box(slide, Inches(1.0), Inches(4.9), Inches(11), Inches(0.3),
             "🚀 未来改进方向", font_size=20, color=PRIMARY, bold=True)

outlook_items = [
    ("框架升级", "迁移至 Spring Boot + MyBatis / Spring Data JPA，提升开发效率"),
    ("前端重构", "采用 Vue.js / React 前后端分离架构，提升用户体验"),
    ("功能增强", "添加图片上传、在线支付（支付宝/微信）、评论评分、秒杀活动"),
    ("性能优化", "引入 Redis 缓存、Druid 连接池、数据库索引优化、CDN 加速"),
]

for i, (tag, desc) in enumerate(outlook_items):
    y = Inches(5.3) + Inches(0.35) * i
    # 标签
    tag_box = add_shape_bg(slide, Inches(1.0), y, Inches(1.5), Inches(0.3), PRIMARY)
    add_text_box(slide, Inches(1.0), y, Inches(1.5), Inches(0.3),
                 tag, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.7), y + Inches(0.02), Inches(9), Inches(0.25),
                 desc, font_size=12, color=RGBColor(0xB2, 0xBE, 0xC3))

add_page_number(slide, 9, 8)


# ============================================================
# 幻灯片 10: 致谢
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PRIMARY)

add_text_box(slide, Inches(0), Inches(2.5), Inches(13.333), Inches(1),
             "感谢聆听", font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.6),
             "THANK YOU", font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)

# 分隔线
add_shape_bg(slide, Inches(5.5), Inches(4.3), Inches(2.333), Inches(0.03), PRIMARY)

add_text_box(slide, Inches(0), Inches(4.6), Inches(13.333), Inches(0.5),
             "简易购物网站 ShoppingWeb", font_size=20, color=RGBColor(0xB2, 0xBE, 0xC3), alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(5.2), Inches(13.333), Inches(0.5),
             "Java 8 + Servlet 4.0 + JSP 2.3 + MySQL 8.0 + Maven + Tomcat 9", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(6.0), Inches(13.333), Inches(0.4),
             "欢迎提问 🙋", font_size=18, color=PRIMARY, alignment=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3), PRIMARY)


# ==================== 保存 ====================
output_path = r"C:\Users\27125\Desktop\xhy的Javaweb作业\ShoppingWeb_答辩PPT.pptx"
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print(f"共 {len(prs.slides)} 张幻灯片")
